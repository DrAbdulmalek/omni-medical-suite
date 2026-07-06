"""
Document Parser Module for Medical Handwriting OCR (HF Spaces Edition).

Simplified version that does not depend on app.config.settings for UPLOAD_DIR.
Uses a local ``./parsed`` directory for output instead.
"""

from __future__ import annotations

import logging
import os
import uuid
from pathlib import Path
from typing import List, Optional

from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

# Use a local directory instead of settings.UPLOAD_DIR
_OUTPUT_DIR = os.environ.get("UPLOAD_DIR", "./parsed")


# =============================================================================
# Pydantic Models
# =============================================================================


class ImageContent(BaseModel):
    page_number: int = Field(..., description="1-based page number where image was found")
    image_path: str = Field(..., description="File path to the extracted image")
    bbox: Optional[dict] = None
    caption: Optional[str] = None
    width: Optional[int] = None
    height: Optional[int] = None
    mime_type: Optional[str] = None


class TableContent(BaseModel):
    page_number: int = Field(...)
    table_index: int = Field(...)
    headers: List[str] = Field(default_factory=list)
    rows: List[List[str]] = Field(default_factory=list)
    confidence: float = Field(default=0.0, ge=0.0, le=1.0)
    bbox: Optional[dict] = None
    row_count: int = Field(0)
    col_count: int = Field(0)


class PageContent(BaseModel):
    page_number: int = Field(...)
    text: str = Field(default="")
    images: List[ImageContent] = Field(default_factory=list)
    tables: List[TableContent] = Field(default_factory=list)
    word_count: int = Field(0)
    has_arabic: bool = Field(False)
    has_latin: bool = Field(False)


class DocumentParseResult(BaseModel):
    document_id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    file_name: str = Field("")
    file_type: str = Field("")
    page_count: int = Field(0)
    pages: List[PageContent] = Field(default_factory=list)
    full_text: str = Field("")
    total_tables: int = Field(0)
    total_images: int = Field(0)
    has_arabic: bool = Field(False)
    processing_time_ms: float = Field(0.0)
    warnings: List[str] = Field(default_factory=list)


# =============================================================================
# Helpers
# =============================================================================


def _has_arabic(text: str) -> bool:
    return any("\u0600" <= c <= "\u06FF" for c in text)


def _has_latin(text: str) -> bool:
    return any(c.isascii() and c.isalpha() for c in text)


def _add_rtl_markers(text: str) -> str:
    if not text:
        return text
    lines: list[str] = []
    for line in text.split("\n"):
        if _has_arabic(line):
            lines.append(f"\u202B{line}\u202C")
        else:
            lines.append(line)
    return "\n".join(lines)


# =============================================================================
# DocumentParser
# =============================================================================


class DocumentParser:

    def __init__(self) -> None:
        os.makedirs(_OUTPUT_DIR, exist_ok=True)
        logger.info("DocumentParser initialized. Output dir: %s", _OUTPUT_DIR)

    def parse_document(
        self,
        file_path: str,
        file_type: Optional[str] = None,
    ) -> DocumentParseResult:
        import time

        start = time.perf_counter()
        file_type = (file_type or "").lower().lstrip(".")
        if not file_type:
            file_type = Path(file_path).suffix.lstrip(".").lower()

        file_name = os.path.basename(file_path)
        logger.info("Parsing document: %s (type=%s)", file_name, file_type)

        result = DocumentParseResult(file_name=file_name, file_type=file_type)

        try:
            if file_type == "pdf":
                result = self._parse_pdf(file_path, result)
            elif file_type == "docx":
                result = self._parse_docx(file_path, result)
            elif file_type in ("pptx", "ppt"):
                result = self._parse_pptx(file_path, result)
            elif file_type in ("html", "htm"):
                result = self._parse_html(file_path, result)
            else:
                result.warnings.append(f"Unsupported file type: {file_type}.")
        except Exception as exc:
            logger.error("Failed to parse %s: %s", file_name, exc, exc_info=True)
            result.warnings.append(f"Parse error: {exc}")

        result.page_count = len(result.pages)
        result.full_text = "\n\n".join(p.text for p in result.pages)
        result.total_tables = sum(len(p.tables) for p in result.pages)
        result.total_images = sum(len(p.images) for p in result.pages)
        result.has_arabic = any(p.has_arabic for p in result.pages)
        result.processing_time_ms = (time.perf_counter() - start) * 1000

        return result

    def _parse_pdf(self, file_path: str, result: DocumentParseResult) -> DocumentParseResult:
        # Try PyMuPDF (fitz)
        try:
            import fitz
            doc = fitz.open(file_path)
            for page_idx in range(len(doc)):
                page = doc.load_page(page_idx)
                text = page.get_text("text") or ""
                result.pages.append(PageContent(
                    page_number=page_idx + 1,
                    text=_add_rtl_markers(text),
                    has_arabic=_has_arabic(text),
                    has_latin=_has_latin(text),
                    word_count=len(text.split()),
                ))
            doc.close()
            return result
        except ImportError:
            pass

        # Try pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    result.pages.append(PageContent(
                        page_number=page_idx + 1,
                        text=_add_rtl_markers(text),
                        has_arabic=_has_arabic(text),
                        has_latin=_has_latin(text),
                        word_count=len(text.split()),
                    ))
            return result
        except ImportError:
            result.warnings.append("Neither PyMuPDF nor pdfplumber available for PDF parsing")
            return result

    def _parse_docx(self, file_path: str, result: DocumentParseResult) -> DocumentParseResult:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            result.warnings.append("python-docx is not installed")
            return result

        try:
            doc = DocxDocument(file_path)
            buffer_lines: List[str] = []
            for para in doc.paragraphs:
                if para.text and para.text.strip():
                    buffer_lines.append(para.text)
            if buffer_lines:
                text = "\n".join(buffer_lines)
                result.pages.append(PageContent(
                    page_number=1,
                    text=_add_rtl_markers(text),
                    has_arabic=_has_arabic(text),
                    has_latin=_has_latin(text),
                    word_count=len(text.split()),
                ))
            return result
        except Exception as exc:
            result.warnings.append(f"DOCX parse error: {exc}")
            return result

    def _parse_pptx(self, file_path: str, result: DocumentParseResult) -> DocumentParseResult:
        try:
            from pptx import Presentation
        except ImportError:
            result.warnings.append("python-pptx is not installed")
            return result

        try:
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text and paragraph.text.strip():
                                texts.append(paragraph.text)
                text = "\n".join(texts)
                result.pages.append(PageContent(
                    page_number=slide_idx + 1,
                    text=_add_rtl_markers(text),
                    has_arabic=_has_arabic(text),
                    has_latin=_has_latin(text),
                    word_count=len(text.split()),
                ))
            return result
        except Exception as exc:
            result.warnings.append(f"PPTX parse error: {exc}")
            return result

    def _parse_html(self, file_path: str, result: DocumentParseResult) -> DocumentParseResult:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            result.warnings.append("beautifulsoup4 is not installed")
            return result

        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                html_content = f.read()
            soup = BeautifulSoup(html_content, "html.parser")
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()
            body = soup.find("body") or soup
            text = body.get_text(separator="\n", strip=True)
            result.pages.append(PageContent(
                page_number=1,
                text=_add_rtl_markers(text),
                has_arabic=_has_arabic(text),
                has_latin=_has_latin(text),
                word_count=len(text.split()),
            ))
            return result
        except Exception as exc:
            result.warnings.append(f"HTML parse error: {exc}")
            return result


# Singleton
document_parser = DocumentParser()
